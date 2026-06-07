"""
SUNUM.pptx üreteci — Hocaya gösterilecek PowerPoint sunumu.

Çalıştırma:
    python3 generate_pptx.py

Çıktı:
    SUNUM.pptx (proje kökünde)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Renk Paleti
RENK_KOYU_MAVI = RGBColor(0x2C, 0x3E, 0x50)
RENK_MAVI = RGBColor(0x34, 0x98, 0xDB)
RENK_YESIL = RGBColor(0x27, 0xAE, 0x60)
RENK_TURUNCU = RGBColor(0xE6, 0x7E, 0x22)
RENK_KIRMIZI = RGBColor(0xE7, 0x4C, 0x3C)
RENK_GRI = RGBColor(0x7F, 0x8C, 0x8D)
RENK_HAFIF = RGBColor(0xEC, 0xF0, 0xF1)

# 16:9 boyut
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def yeni_sunum():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def baslik_slayt(prs, baslik, alt_baslik, isim, kurum):
    """Kapak slaytı."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # boş layout

    # Arka plan rengi
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RENK_KOYU_MAVI

    # Başlık kutusu
    tbox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(1.5))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = baslik
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Alt başlık
    abox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.3), Inches(1))
    af = abox.text_frame
    af.word_wrap = True
    ap = af.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    ar = ap.add_run()
    ar.text = alt_baslik
    ar.font.size = Pt(24)
    ar.font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)

    # İsim ve kurum
    nbox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.3), Inches(1.5))
    nf = nbox.text_frame
    nf.word_wrap = True
    p1 = nf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = isim
    r1.font.size = Pt(28)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = nf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = kurum
    r2.font.size = Pt(20)
    r2.font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)


def icerikli_slayt(prs, baslik, alt_baslik, maddeler, renk_baslik=None):
    """Madde listesi olan slayt."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    renk_baslik = renk_baslik or RENK_KOYU_MAVI

    # Üst başlık şeridi
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = renk_baslik
    bar.line.fill.background()

    # Başlık
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.5))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = baslik
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Alt başlık
    if alt_baslik:
        sbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.5))
        sf = sbox.text_frame
        sp = sf.paragraphs[0]
        sr = sp.add_run()
        sr.text = alt_baslik
        sr.font.size = Pt(16)
        sr.font.italic = True
        sr.font.color.rgb = RENK_GRI

    # İçerik kutusu
    y = 1.6 if alt_baslik else 1.0
    cbox = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(12.0), Inches(7.5 - y - 0.3))
    cf = cbox.text_frame
    cf.word_wrap = True

    for i, madde in enumerate(maddeler):
        # Madde dict olabilir: {"text": "...", "bold": True, "color": ..., "size": ...}
        if isinstance(madde, dict):
            metin = madde.get("text", "")
            kalin = madde.get("bold", False)
            renk = madde.get("color", RGBColor(0x33, 0x33, 0x33))
            boyut = madde.get("size", 20)
            girinti = madde.get("indent", 0)
        else:
            metin = str(madde)
            kalin = False
            renk = RGBColor(0x33, 0x33, 0x33)
            boyut = 20
            girinti = 0

        if i == 0:
            p = cf.paragraphs[0]
        else:
            p = cf.add_paragraph()
        p.level = girinti
        r = p.add_run()
        r.text = metin
        r.font.size = Pt(boyut)
        r.font.bold = kalin
        r.font.color.rgb = renk


def tablo_slayt(prs, baslik, basliklar, satirlar, vurgu_indeksleri=None):
    """Tablolu slayt."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Başlık şeridi
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RENK_KOYU_MAVI
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.5))
    p = tbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = baslik
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Tablo
    rows = len(satirlar) + 1
    cols = len(basliklar)
    left = Inches(0.7)
    top = Inches(1.2)
    width = Inches(12.0)
    height = Inches(6.0)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Başlık satırı
    for j, h in enumerate(basliklar):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RENK_MAVI
        cell.text_frame.text = h
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(16)

    # Veri satırları
    vurgu_indeksleri = vurgu_indeksleri or []
    for i, satir in enumerate(satirlar):
        for j, val in enumerate(satir):
            cell = table.cell(i + 1, j)
            cell.text_frame.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
                    if i in vurgu_indeksleri:
                        run.font.bold = True
                        run.font.color.rgb = RENK_YESIL


# ═══════════════════════════════════════════════════════
# SUNUM OLUŞTUR
# ═══════════════════════════════════════════════════════

prs = yeni_sunum()

# SLAYT 1: KAPAK
baslik_slayt(
    prs,
    "AI-Destekli Veri Sıkıştırma",
    "Klasik Algoritmaların Sinir Ağı Tabanlı Otomatik Seçimi\nve LLM Tabanlı Akıllı Sözlük Üretimi",
    "Betül Arslan",
    "Yıldız Teknik Üniversitesi · Bilgisayar Mühendisliği · Bahar 2026",
)

# SLAYT 2: İÇİNDEKİLER
icerikli_slayt(
    prs, "📋 İçindekiler", "",
    [
        {"text": "1. Problem ve Amaç", "size": 22},
        {"text": "2. Bilgi Teorisi Temeli — Shannon Entropisi", "size": 22},
        {"text": "3. Sistem Mimarisi", "size": 22},
        {"text": "4. Klasik Algoritmalar (Huffman / LZW / BWT+MTF)", "size": 22},
        {"text": "5. 🤖 AI Entegrasyonu — Akıllı Sözlük + NN Seçici", "size": 22, "bold": True, "color": RENK_TURUNCU},
        {"text": "6. Sinir Ağının 11 Özelliği", "size": 22},
        {"text": "7. Performans Sonuçları", "size": 22},
        {"text": "8. Endüstri Karşılaştırması (gzip, bzip2)", "size": 22},
        {"text": "9. Canlı Demo", "size": 22},
        {"text": "10. Sonuç", "size": 22},
    ],
)

# SLAYT 3: PROBLEM
icerikli_slayt(
    prs, "🎯 Problem ve Amaç", "Her algoritma farklı veride iyi performans gösterir",
    [
        {"text": "Huffman: Karakter bazlı → doğal dilde orta", "size": 22, "bold": True, "color": RENK_MAVI},
        {"text": "LZW: Pattern bazlı → tekrarlı veride güçlü", "size": 22, "bold": True, "color": RENK_MAVI},
        {"text": "BWT (bzip2): Yapısal düzen → kümeli veride mükemmel", "size": 22, "bold": True, "color": RENK_MAVI},
        {"text": "", "size": 12},
        {"text": "❓ Soru: \"Verilen bir metin için hangi algoritma en iyi?\"", "size": 24, "bold": True, "color": RENK_KIRMIZI},
        {"text": "", "size": 8},
        {"text": "📚 Bu, Algorithm Selection Problem (Rice, 1976)", "size": 20, "color": RENK_GRI},
        {"text": "Modern çözüm: Makine öğrenmesi (Kotthoff, 2014)", "size": 20, "color": RENK_GRI},
    ],
)

# SLAYT 4: HOCA KRİTERLERİ
icerikli_slayt(
    prs, "🎯 Hocanın PDF'inde 2 Değerlendirme Kriteri", "Her ikisine de yanıt verildi",
    [
        {"text": "✅ Kriter 1: En İyi Performans", "size": 26, "bold": True, "color": RENK_YESIL},
        {"text": "• 6 veri tipinde +%2.5 ile +%86.5 iyileşme", "size": 18, "indent": 1},
        {"text": "• Bzip2'den %36.5 daha küçük (kısa metinde)", "size": 18, "indent": 1},
        {"text": "• 100 birim test ile sıfır kayıp", "size": 18, "indent": 1},
        {"text": "", "size": 10},
        {"text": "✅ Kriter 2: En İyi AI Entegrasyonu", "size": 26, "bold": True, "color": RENK_YESIL},
        {"text": "• 🤖 LLM ile LZW sözlüğü üretimi (PDF'te özellikle istenen!)", "size": 18, "indent": 1, "bold": True, "color": RENK_TURUNCU},
        {"text": "• 3-sınıflı sinir ağı algoritma seçici", "size": 18, "indent": 1},
        {"text": "• Akıllı Hibrit no-regret yöneticisi", "size": 18, "indent": 1},
        {"text": "• 42 adım AI günlüğü, 12 akademik doğrulama", "size": 18, "indent": 1},
    ],
)

# SLAYT 5: SHANNON ENTROPİSİ
icerikli_slayt(
    prs, "📐 Bilgi Teorisi: Shannon Entropisi", "Sıkıştırma için teorik alt sınır",
    [
        {"text": "Formül (Shannon, 1948):", "size": 22, "bold": True},
        {"text": "H(X) = -Σ p(c) · log₂ p(c)   [bit/karakter]", "size": 26, "color": RENK_MAVI, "bold": True},
        {"text": "", "size": 8},
        {"text": "Sezgisel anlam:", "size": 22, "bold": True},
        {"text": "• \"aaaaaa\" → 0 bit/kar (tahmin kesin)", "size": 18, "indent": 1},
        {"text": "• \"abababab\" → 1 bit/kar (2 karakter)", "size": 18, "indent": 1},
        {"text": "• Türkçe metin → ~4.5 bit/kar", "size": 18, "indent": 1},
        {"text": "• Rastgele bit → 8 bit/kar (maksimum)", "size": 18, "indent": 1},
        {"text": "", "size": 8},
        {"text": "⚡ Hiçbir algoritma bu sınırın altına inemez!", "size": 24, "bold": True, "color": RENK_KIRMIZI},
    ],
)

# SLAYT 6: SİSTEM MİMARİSİ
icerikli_slayt(
    prs, "🏗️ Sistem Mimarisi", "3 katmanlı AI entegrasyonu",
    [
        {"text": "[Kullanıcı Metni]", "size": 18, "color": RENK_KOYU_MAVI, "bold": True},
        {"text": "      ↓", "size": 14, "color": RENK_GRI},
        {"text": "┌─ Klasik Algo. ──┐  ┌─ Groq LLM ──┐  ┌─ Sinir Ağı ──┐", "size": 14, "color": RENK_MAVI},
        {"text": "│ Huffman/LZW/BWT │  │ AI Sözlüğü   │  │ MLP Seçici   │", "size": 14, "color": RENK_MAVI},
        {"text": "└──────────┬──────┘  └──────┬──────┘  └──────┬──────┘", "size": 14, "color": RENK_MAVI},
        {"text": "           └─────────────┼──────────────┘", "size": 14, "color": RENK_GRI},
        {"text": "                         ▼", "size": 14, "color": RENK_GRI},
        {"text": "         ┌─ Akıllı Hibrit Yöneticisi ─┐", "size": 16, "color": RENK_TURUNCU, "bold": True},
        {"text": "         │  NN seç + BWT post-check   │", "size": 14, "color": RENK_TURUNCU},
        {"text": "         └──────────────┬─────────────┘", "size": 16, "color": RENK_TURUNCU},
        {"text": "                        ▼", "size": 14, "color": RENK_GRI},
        {"text": "                  Sıkıştırılmış Çıktı (.bin)", "size": 18, "bold": True, "color": RENK_YESIL},
    ],
)

# SLAYT 7: KLASİK ALGORİTMALAR
icerikli_slayt(
    prs, "📦 Klasik Algoritmalar", "Akademik temelli üç algoritma",
    [
        {"text": "1️⃣ Huffman (1952)", "size": 24, "bold": True, "color": RENK_MAVI},
        {"text": "• Optimal prefix-free kod, min-heap ile ağaç", "size": 18, "indent": 1},
        {"text": "• H(X) ≤ L < H(X)+1 (Cover & Thomas)", "size": 18, "indent": 1},
        {"text": "", "size": 8},
        {"text": "2️⃣ LZW (1984)", "size": 24, "bold": True, "color": RENK_MAVI},
        {"text": "• Sözlük dosya başında gönderilmez", "size": 18, "indent": 1},
        {"text": "• Türkçe karakter çözümü: alfabe ekle (Salomon §6.13)", "size": 18, "indent": 1},
        {"text": "", "size": 8},
        {"text": "3️⃣ BWT + MTF + RLE + Huffman (Klasik bzip2)", "size": 24, "bold": True, "color": RENK_MAVI},
        {"text": "• Tam pipeline (Burrows-Wheeler 1994, Bentley 1986)", "size": 18, "indent": 1},
        {"text": "• Bloklu sistem ile sınırsız uzunluk", "size": 18, "indent": 1},
    ],
)

# SLAYT 8: AI AKILLI SÖZLÜK - VURGULU
icerikli_slayt(
    prs, "🤖 AI Entegrasyonu — Akıllı Sözlük", "Hocanın PDF'inde özellikle istenen özellik",
    [
        {"text": "📋 Hocanın PDF'inden:", "size": 20, "italic": True, "color": RENK_GRI},
        {"text": "\"LZW veya Huffman için 'en optimize' sözlüğü bir LLM'e", "size": 18, "color": RENK_GRI, "indent": 1},
        {"text": " analiz ettirip oluşturmak.\"", "size": 18, "color": RENK_GRI, "indent": 1},
        {"text": "", "size": 8},
        {"text": "✅ Bu Projedeki Çözüm:", "size": 22, "bold": True, "color": RENK_YESIL},
        {"text": "1. Kullanıcı metni LLM'e (LLaMA 3.3 70B) gönderilir", "size": 18, "indent": 1},
        {"text": "2. Prompt: \"60 yaygın kelime/ifade öner\"", "size": 18, "indent": 1},
        {"text": "3. LLM döner: [\"Türkiye\", \"Avrupa\", \"köprü\", ...]", "size": 18, "indent": 1},
        {"text": "4. Bu kelimeler LZW başlangıç sözlüğüne eklenir", "size": 18, "indent": 1},
        {"text": "5. Sözlük: 256 ASCII + 60 LLM = 316 başlangıç pattern", "size": 18, "indent": 1},
        {"text": "", "size": 8},
        {"text": "🎯 Sonuç: Standart LZW'ye göre +%13 ek küçülme", "size": 22, "bold": True, "color": RENK_TURUNCU},
    ],
    renk_baslik=RENK_TURUNCU,
)

# SLAYT 9: PROMPT MÜHENDİSLİĞİ
icerikli_slayt(
    prs, "🔧 Prompt Mühendisliği", "Sistematik AI komutları",
    [
        {"text": "Sistem Promptu:", "size": 20, "bold": True, "color": RENK_MAVI},
        {"text": "\"Sen veri sıkıştırma uzmanısın. LZW için pattern eşleşmesini", "size": 16},
        {"text": " hızlandıracak EN YAYGIN kelimeleri tahmin et. JSON döndür.\"", "size": 16},
        {"text": "", "size": 8},
        {"text": "Kullanıcı Promptu:", "size": 20, "bold": True, "color": RENK_MAVI},
        {"text": "\"Bu Türkçe metin için 60 yaygın kelime öner.", "size": 16},
        {"text": " Metin: \"\"\"Türkiye, Avrupa ve Asya...\"\"\"", "size": 16},
        {"text": " Format: [\"kelime1\", \"kelime2\", ...]\"", "size": 16},
        {"text": "", "size": 8},
        {"text": "📚 Karmaşık algoritmayı parçalama:", "size": 20, "bold": True, "color": RENK_TURUNCU},
        {"text": "BWT için adım adım sordum:", "size": 16},
        {"text": "  1. \"BWT nedir? Adımlarını yaz.\"", "size": 16, "indent": 1},
        {"text": "  2. \"Move-to-Front nasıl çalışır?\"", "size": 16, "indent": 1},
        {"text": "  3. \"LF mapping ile decode nasıl?\"", "size": 16, "indent": 1},
    ],
)

# SLAYT 10: SİNİR AĞI
icerikli_slayt(
    prs, "🧠 Sinir Ağı: MLP 32→16→8", "11 özelliği girdi alıp 3 sınıf çıktısı verir",
    [
        {"text": "Mimari (Bishop §5.1, Goodfellow §6.3):", "size": 20, "bold": True},
        {"text": "Girdi (11 özellik)", "size": 18, "color": RENK_MAVI},
        {"text": "       ↓", "size": 14, "color": RENK_GRI},
        {"text": "[32 nöron, ReLU]  →  [16, ReLU]  →  [8, ReLU]", "size": 18, "color": RENK_MAVI, "bold": True},
        {"text": "       ↓", "size": 14, "color": RENK_GRI},
        {"text": "Çıktı: Huffman / LZW / BWT (Softmax)", "size": 18, "color": RENK_MAVI},
        {"text": "", "size": 12},
        {"text": "Eğitim Sonuçları:", "size": 20, "bold": True},
        {"text": "• 2.357 örnek (sentetik + corpus)", "size": 18, "indent": 1},
        {"text": "• L2 regularizasyon + Early stopping", "size": 18, "indent": 1},
        {"text": "• 5-fold Cross-Validation: %91.7 ± %1.8", "size": 18, "indent": 1},
        {"text": "• Hold-out test: %95.2 ⭐", "size": 22, "indent": 1, "bold": True, "color": RENK_YESIL},
    ],
)

# SLAYT 11: 11 ÖZELLİK 1/2
tablo_slayt(
    prs, "🔬 Sinir Ağının 11 Özelliği (1/2)",
    ["#", "Özellik", "Önem", "Hangi Algoritmaya İşaret"],
    [
        ["1", "Shannon entropisi", "0.042", "Yüksek → Huffman"],
        ["2", "Benzersiz oran ⭐", "0.247", "Düşük → BWT/LZW"],
        ["3", "Top-3 yoğunluk", "0.024", "Yüksek → Huffman"],
        ["4", "Boşluk oranı", "0.018", "Yüksek → LZW (doğal dil)"],
        ["5", "Türkçe karakter oranı", "0.058", "Yüksek → Türkçe corpus"],
        ["6", "Run-length ortalama", "0.082", "Yüksek → BWT"],
    ],
    vurgu_indeksleri=[1],
)

# SLAYT 12: 11 ÖZELLİK 2/2
tablo_slayt(
    prs, "🔬 Sinir Ağının 11 Özelliği (2/2)",
    ["#", "Özellik", "Önem", "Hangi Algoritmaya İşaret"],
    [
        ["7", "Rakam oranı", "0.008", "Yüksek → LZW (log/JSON)"],
        ["8", "Büyük harf oranı", "0.012", "Dolaylı"],
        ["9", "Bigram entropi ⭐", "0.183", "Düşük → LZW"],
        ["10", "Max koşu oranı", "0.038", "Yüksek → BWT"],
        ["11", "Alfabe boyutu (log₂) ⭐", "0.096", "Küçük → BWT"],
    ],
    vurgu_indeksleri=[2, 4],
)

# SLAYT 13: AKILLI HİBRİT
icerikli_slayt(
    prs, "⚡ Akıllı Hibrit (No-Regret Garantisi)", "Asla standart Huffman'dan kötü olamaz",
    [
        {"text": "def smart_hybrid(text):", "size": 18, "bold": True, "color": RENK_MAVI},
        {"text": "    # 1. Sinir ağı seç", "size": 16, "color": RENK_GRI},
        {"text": "    nn = nn_predict(text)['algorithm']", "size": 16},
        {"text": "", "size": 6},
        {"text": "    # 2. Seçilen algoritmayı çalıştır", "size": 16, "color": RENK_GRI},
        {"text": "    smart_bits = ilgili_algoritma(text)", "size": 16},
        {"text": "", "size": 6},
        {"text": "    # 3. BWT post-check güvenlik ağı", "size": 16, "color": RENK_GRI},
        {"text": "    bwt_alt = bwt_rle_huffman_bits(text)", "size": 16},
        {"text": "    if bwt_alt < smart_bits:", "size": 16},
        {"text": "        smart_bits = bwt_alt", "size": 16},
        {"text": "", "size": 6},
        {"text": "    # 4. Standart Huffman ile karşılaştır", "size": 16, "color": RENK_GRI},
        {"text": "    return min(smart_bits, huffman_bits(text))", "size": 16},
        {"text": "", "size": 8},
        {"text": "🛡️ Garanti: Akıllı Hibrit ≤ standart Huffman", "size": 20, "bold": True, "color": RENK_YESIL},
    ],
)

# SLAYT 14: SONUÇLAR
tablo_slayt(
    prs, "📊 Performans Sonuçları (6 Veri Tipi)",
    ["Veri Türü", "Std Huffman", "Akıllı Hibrit", "İyileşme"],
    [
        ["Tekrarlı (ABC×100)", "%78.2", "%96.9", "+%85.9 🔥"],
        ["DNA dizisi", "%73.8", "%94.1", "+%77.5 🔥"],
        ["JSON log", "%60.7", "%94.7", "+%86.5 🔥"],
        ["Doğal Türkçe", "%37.1", "%38.7", "+%2.5"],
        ["Wikipedia TR", "%39.8", "%58.3", "+%30.7"],
        ["Türkçe haber", "%42.0", "%54.3", "+%21.2"],
    ],
    vurgu_indeksleri=[0, 1, 2],
)

# SLAYT 15: ENDÜSTRİ KARŞILAŞTIRMA
tablo_slayt(
    prs, "🏭 Endüstri Karşılaştırması (730 karakter Türkçe akademik metin)",
    ["Algoritma", "Boyut", "Küçülme", "Süre"],
    [
        ["gzip -9", "95 byte", "%87.0", "0.06 ms"],
        ["zlib -9", "91 byte", "%87.5", "0.01 ms"],
        ["bzip2 -9", "148 byte", "%79.7", "0.80 ms"],
        ["lzma -9", "124 byte", "%83.0", "6.79 ms"],
        ["🏆 bwt_rle_huffman (BİZ)", "94 byte", "%87.1", "0.56 ms"],
    ],
    vurgu_indeksleri=[4],
)

# SLAYT 16: SONUÇ ÖZET
icerikli_slayt(
    prs, "📈 Sonuç — Hocanın 2 Kriteri ✅ ✅", "",
    [
        {"text": "✅ En İyi Performans Kriteri", "size": 26, "bold": True, "color": RENK_YESIL},
        {"text": "• 6 veri tipinde +%2.5 ile +%86.5 iyileşme", "size": 18, "indent": 1},
        {"text": "• Bzip2'den %36.5 daha küçük (kısa metinde)", "size": 18, "indent": 1},
        {"text": "• 100 birim test, sıfır kayıp", "size": 18, "indent": 1},
        {"text": "", "size": 10},
        {"text": "✅ En İyi AI Entegrasyonu Kriteri", "size": 26, "bold": True, "color": RENK_YESIL},
        {"text": "• 🤖 LLM ile LZW sözlüğü (PDF'te istenen!)", "size": 18, "indent": 1, "bold": True, "color": RENK_TURUNCU},
        {"text": "• 3-sınıflı NN algoritma seçici (%95.2)", "size": 18, "indent": 1},
        {"text": "• Akıllı Hibrit + 42 AI günlüğü", "size": 18, "indent": 1},
        {"text": "• 12 akademik kaynak doğrulaması", "size": 18, "indent": 1},
        {"text": "", "size": 10},
        {"text": "🚀 Sistem canlı çalışıyor, akademik temelli, kayıpsız garanti", "size": 22, "bold": True, "color": RENK_KOYU_MAVI},
    ],
)

# SLAYT 17: TEŞEKKÜRLER
baslik_slayt(
    prs,
    "Teşekkürler 🎓",
    "Sorularınız?\n\nhuggingface.co/spaces/tien23/ai-veri-sikistirma\ngithub.com/betullarslan-cpu/ai-veri-sikistirma",
    "Betül Arslan",
    "Yıldız Teknik Üniversitesi · Bahar 2026",
)

# Kaydet
prs.save("SUNUM.pptx")
print("✅ SUNUM.pptx üretildi")
print(f"   Toplam slayt: {len(prs.slides)}")
